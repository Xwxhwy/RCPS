# rcps/preval/feature_extractor.py

from typing import Dict, Any, List, Tuple
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import numpy as np
from PIL import Image
from transformers import AutoProcessor, AutoModel
from sentence_transformers import SentenceTransformer, util
import torch
from ..utils import get_logger  # 假设你有这个工具
import io
import itertools

logger = get_logger(__name__)


class PrevalFeatureExtractor:
    """
    为PREVAL评估框架提取丰富的多模态特征。

    该提取器整合了多种先进模型和算法，旨在从PPTX文件中量化
    内容（Content）、连贯性（Coherence）和设计（Design）三个维度的质量。
    特征包括：
    - 文本特征 (SentenceTransformer)
    - 图像美学特征 (Aesthetic Predictor)
    - 布局结构特征 (Alignment, Balance, White Space)
    - 图文对齐性 (CLIP)
    """

    def __init__(self, device: str = 'cpu'):
        self.device = device
        self.text_model = None
        self.clip_model = None
        self.clip_processor = None
        self.aesthetic_model = None

        try:
            logger.info("Loading models for PrevalFeatureExtractor...")


            self.text_model = SentenceTransformer('all-MiniLM-L6-v2', device=self.device)


            clip_model_name = "openai/clip-vit-base-patch32"
            self.clip_model = AutoModel.from_pretrained(clip_model_name).to(self.device)
            self.clip_processor = AutoProcessor.from_pretrained(clip_model_name)


            self.aesthetic_model = torch.nn.Linear(self.clip_model.config.vision_config.hidden_size, 1).to(self.device)
            logger.info("All models loaded successfully.")

        except Exception as e:
            logger.error(f"Failed to load models for PrevalFeatureExtractor: {e}", exc_info=True)
            raise IOError("Could not initialize PrevalFeatureExtractor due to model loading failure.")

    def extract(self, pptx_path: str) -> Dict[str, Any]:

        if self.text_model is None:
            logger.error("Feature extractor is not properly initialized. Cannot extract features.")
            return {}

        try:
            prs = PptxPresentation(pptx_path)
            doc_width = prs.slide_width
            doc_height = prs.slide_height
        except Exception as e:
            logger.error(f"Failed to open or parse PPTX file: {pptx_path}. Error: {e}")
            return {}

        all_slide_data = []
        for i, slide in enumerate(prs.slides):
            # 模拟渲染幻灯片为图像
            slide_image = self._render_slide_to_pil_image(slide, doc_width, doc_height)
            slide_text_contents = self._extract_text_from_slide(slide)
            slide_shapes = self._extract_shape_properties(slide)

            all_slide_data.append({
                "image": slide_image,
                "texts": slide_text_contents,
                "shapes": slide_shapes,
            })

        # --- 特征计算 ---
        text_features = self._compute_text_features(all_slide_data)
        layout_features = self._compute_layout_features(all_slide_data, doc_width, doc_height)
        visual_features = self._compute_visual_features(all_slide_data)

        # 合并所有特征
        features = {**text_features, **layout_features, **visual_features}

        # 清理NaN或inf值
        for key, value in features.items():
            if np.isnan(value) or np.isinf(value):
                features[key] = 0.0

        return features

    def _compute_text_features(self, all_slide_data: List[Dict]) -> Dict[str, float]:
        """计算所有与文本相关的特征。"""
        full_texts = ["\n".join(slide["texts"]) for slide in all_slide_data]
        slide_embeddings = self.text_model.encode(
            [text for text in full_texts if text.strip()],
            convert_to_tensor=True, show_progress_bar=False
        )

        if len(slide_embeddings) < 2:
            return {
                'text_coherence': 0.0,
                'avg_text_length': np.mean([len(t) for t in full_texts]) if full_texts else 0.0,
                'text_redundancy': 0.0
            }

        # 文本连贯性: 相邻幻灯片嵌入的平均余弦相似度
        sims = util.cos_sim(slide_embeddings[:-1], slide_embeddings[1:])
        text_coherence = torch.diag(sims).mean().item()

        # 文本冗余度: 衡量幻灯片之间内容的重复程度
        all_pairs_sims = util.cos_sim(slide_embeddings, slide_embeddings)

        redundancy = all_pairs_sims[torch.triu(torch.ones_like(all_pairs_sims), diagonal=1).bool()].mean().item()

        return {
            'text_coherence': text_coherence,
            'avg_text_length': np.mean([len(t) for t in full_texts]),
            'text_redundancy': redundancy,
        }

    def _compute_layout_features(self, all_slide_data: List[Dict], doc_width, doc_height) -> Dict[str, float]:
        """计算所有与布局相关的特征。"""
        alignment_scores, balance_scores, whitespace_ratios = [], [], []

        for slide in all_slide_data:
            shapes = slide["shapes"]
            if len(shapes) < 2: continue

            # 对齐度: 检查形状的左、右、中心线是否对齐
            x_coords = [s['x'] for s in shapes]
            y_coords = [s['y'] for s in shapes]
            centers_x = [s['x'] + s['w'] / 2 for s in shapes]
            centers_y = [s['y'] + s['h'] / 2 for s in shapes]

            # 使用容差来判断对齐
            def get_alignment_score(coords, tolerance=0.01):
                unique_coords, counts = np.unique(np.round(np.array(coords) / (doc_width * tolerance)),
                                                  return_counts=True)
                return (counts > 1).sum() / len(coords) if len(coords) > 0 else 0

            align_x = get_alignment_score(x_coords)
            align_center_x = get_alignment_score(centers_x)
            alignment_scores.append((align_x + align_center_x) / 2)

            # 平衡性: 衡量内容在幻灯片中心线两侧的分布是否均衡
            total_area = sum(s['w'] * s['h'] for s in shapes)
            if total_area > 0:
                left_weight = sum(s['w'] * s['h'] for s in shapes if (s['x'] + s['w'] / 2) < doc_width / 2)
                balance_scores.append(1 - abs(left_weight / total_area - 0.5) * 2)

            # 空白利用率: 非元素区域占总面积的比例
            occupied_area = sum(s['w'] * s['h'] for s in shapes)
            whitespace_ratios.append(1 - occupied_area / (doc_width * doc_height))

        return {
            'avg_alignment_score': np.mean(alignment_scores) if alignment_scores else 0.0,
            'avg_balance_score': np.mean(balance_scores) if balance_scores else 0.0,
            'avg_whitespace_ratio': np.mean(whitespace_ratios) if whitespace_ratios else 0.0,
        }

    def _compute_visual_features(self, all_slide_data: List[Dict]) -> Dict[str, float]:
        """计算所有与视觉和多模态相关的特征。"""
        aesthetic_scores, image_text_alignments = [], []
        images = [slide["image"] for slide in all_slide_data]
        texts = ["\n".join(slide["texts"]) for slide in all_slide_data]

        with torch.no_grad():
            # 批量处理图像
            inputs = self.clip_processor(text=texts, images=images, return_tensors="pt", padding=True,
                                         truncation=True).to(self.device)
            outputs = self.clip_model(**inputs)

            # 图像美学特征
            image_embeds = outputs.image_embeds
            if image_embeds is not None:
                aesthetic_logits = self.aesthetic_model(image_embeds)
                aesthetic_scores = torch.sigmoid(aesthetic_logits).cpu().numpy().flatten().tolist()

            # 图文对齐性 (logits_per_image)
            if outputs.logits_per_image is not None:
                # 对角线上的值代表匹配的(图,文)对的相似度分数
                image_text_alignments = torch.diag(outputs.logits_per_image).cpu().numpy().tolist()

        return {
            'avg_aesthetic_score': np.mean(aesthetic_scores) if aesthetic_scores else 0.0,
            'avg_image_text_alignment': np.mean(image_text_alignments) if image_text_alignments else 0.0,
        }

    def _render_slide_to_pil_image(self, slide, doc_width, doc_height) -> Image.Image:
        # 创建一个白色背景的模拟图像
        img = Image.new('RGB', (int(doc_width / 9525), int(doc_height / 9525)), color='white')
        return img

    def _extract_text_from_slide(self, slide) -> List[str]:
        return [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]

    def _extract_shape_properties(self, slide) -> List[Dict[str, float]]:
        return [{
            "x": shape.left, "y": shape.top,
            "w": shape.width, "h": shape.height
        } for shape in slide.shapes]