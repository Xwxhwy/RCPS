# rcps/renderer.py (修正版)
import os
from abc import ABC, abstractmethod
from typing import Dict

from pptx import Presentation as PptxPresentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .presentation import Presentation
from .shapes import Shape, TextElement, ImageElement
from .utils import get_logger, ensure_dir
from .exceptions import GenerationError
from . import constants as C

logger = get_logger(__name__)


class BaseRenderer(ABC):
    """渲染器的抽象基类。"""

    @abstractmethod
    def render(self, presentation: Presentation, output_path: str):
        pass


class PptxRenderer(BaseRenderer):
    """将 RCPS Presentation 对象渲染为 .pptx 文件。"""

    def render(self, presentation: Presentation, output_path: str):
        """执行渲染过程。"""
        logger.info(f"Rendering presentation '{presentation.presentation_id}' to '{output_path}'...")
        ensure_dir(os.path.dirname(output_path))

        try:
            prs = PptxPresentation()
            prs.slide_width = Pt(720)
            prs.slide_height = Pt(405)

            style_config = self._get_style_config(presentation.design_style, presentation.color_palette)

            for slide_page in presentation.slides:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)

                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(style_config['background_color'][1:])

                for shape in sorted(slide_page.shapes, key=lambda s: s.z_order):
                    self._render_shape(slide, shape, style_config)

            prs.save(output_path)
            logger.info("Presentation rendered successfully.")

        except Exception as e:
            logger.error(f"Failed to render presentation to PPTX: {e}", exc_info=True)
            raise GenerationError("PPTX rendering failed.") from e

    def _render_shape(self, slide: "pptx.slide.Slide", shape: Shape, style_config: Dict):
        """渲染单个Shape到幻灯片上。"""
        try:
            if isinstance(shape, TextElement):
                self._render_text_box(slide, shape, style_config)
            elif isinstance(shape, ImageElement):
                self._render_image(slide, shape)
            else:
                logger.warning(f"Unsupported shape type '{shape.shape_type}' for rendering. Skipping.")
        except Exception as e:
            logger.error(f"Failed to render shape {shape.element_id}: {e}")

    def _render_text_box(self, slide, shape: TextElement, style_config: Dict):
        txBox = slide.shapes.add_textbox(
            Pt(shape.left), Pt(shape.top), Pt(shape.width), Pt(shape.height)
        )
        tf = txBox.text_frame
        tf.text = shape.content
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.alignment = style_config.get('text_align', PP_ALIGN.LEFT)
        font = para.font
        font.name = style_config.get('font_family', 'Calibri')

        if "title" in shape.element_id.lower():
            font.size = Pt(style_config.get('font_size_title', 32))
            font.bold = True
            font.color.rgb = RGBColor.from_string(style_config['primary_color'][1:])
        else:
            font.size = Pt(style_config.get('font_size_body', 18))
            font.color.rgb = RGBColor.from_string(style_config['text_color'][1:])

    def _render_image(self, slide, shape: ImageElement):
        if not shape.content or not os.path.exists(shape.content):
            logger.warning(f"Image path not found or empty for shape {shape.element_id}. Skipping.")
            return

        slide.shapes.add_picture(
            shape.content, Pt(shape.left), Pt(shape.top), width=Pt(shape.width)
        )

    def _get_style_config(self, design_style: str, color_palette: Dict) -> Dict:
        """根据设计风格和色板返回具体的样式配置。"""
        # 修复了对 C 的引用
        config = {
            'font_family': 'Helvetica Neue',
            'font_size_title': 36,
            'font_size_body': 16,
            'text_align': PP_ALIGN.LEFT,
            'primary_color': color_palette['primary'],
            'text_color': color_palette['text'],
            'background_color': color_palette['background']
        }
        if design_style == C.DESIGN_STYLE_MINIMALIST:
            config['font_family'] = 'Roboto'
            config['font_size_title'] = 40
        return config