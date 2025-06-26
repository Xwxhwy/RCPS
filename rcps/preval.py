# rcps/preval.py
from typing import Dict
import torch
from pptx import Presentation as PptxPresentation
from .utils import get_logger

logger = get_logger(__name__)


class PREVAL:
    """
    Preference-based Evaluation Framework via Learned Assessment.
    """

    def __init__(self, model_path: str):
        try:
            self.model = "mock_preval_model"
            logger.info("PREVAL running in MOCK mode.")
            # ------------
        except Exception as e:
            logger.error(f"Failed to load PREVAL model from {model_path}: {e}")
            raise

    def _extract_features(self, prs: PptxPresentation) -> Dict:
        features = {}
        text_content = [shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame]
        features['avg_text_per_slide'] = len("".join(text_content)) / len(prs.slides) if prs.slides else 0

        features['num_slides'] = len(prs.slides)
        features['num_images'] = sum(
            1 for slide in prs.slides for shape in slide.shapes if "Picture" in str(type(shape)))


        logger.info(f"Extracted features: {features}")
        return features

    def evaluate(self, pptx_path: str) -> Dict[str, float]:
        try:
            prs = PptxPresentation(pptx_path)
            features = self._extract_features(prs)

            if self.model == "mock_preval_model":
                # --- 模拟评估逻辑 ---
                score_content = 0.85 - features['avg_text_per_slide'] / 1000
                score_design = 0.7 + features['num_images'] * 0.05
                score_coherence = 0.8
                overall = (score_content + score_design + score_coherence) / 3

                scores = {
                    "Content": round(score_content, 2),
                    "Design": round(score_design, 2),
                    "Coherence": round(score_coherence, 2),
                    "Overall": round(overall, 2),
                }
                logger.info(f"PREVAL mock scores for '{pptx_path}': {scores}")
                return scores

        except Exception as e:
            logger.error(f"Failed to evaluate presentation '{pptx_path}': {e}")
            return {}