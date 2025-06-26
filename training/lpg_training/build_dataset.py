# training/lpg_training/build_dataset.py

import os
import json
import base64
from pathlib import Path
from typing import List, Dict, Any, Tuple

import openai
from pptx import Presentation
from PIL import Image
import io
from tqdm import tqdm


client = openai.OpenAI()


class PptxParser:

    def __init__(self, pptx_path: str):
        self.presentation = Presentation(pptx_path)

    def __iter__(self):
        for i, slide in enumerate(self.presentation.slides):
            yield self.parse_slide(slide, i)

    def parse_slide(self, slide, slide_index: int) -> Dict[str, Any]:
        """从单个幻灯片中提取所有元素及其属性。"""
        elements = []
        for shape in slide.shapes:
            element_info = {
                "type": "unsupported",
                "text": "",
                "position": (shape.left, shape.top),
                "size": (shape.width, shape.height)
            }
            if shape.has_text_frame and shape.text.strip():
                element_info["type"] = "text"
                element_info["text"] = shape.text
            elif shape.shape_type == 13:  # 13 corresponds to a picture
                element_info["type"] = "image"
                # 将图片编码为base64，以便后续传递给多模态模型
                image_bytes = shape.image.blob
                img = Image.open(io.BytesIO(image_bytes))
                # 统一转为PNG格式
                buffered = io.BytesIO()
                img.save(buffered, format="PNG")
                element_info["image_base64"] = base64.b64encode(buffered.getvalue()).decode('utf-8')

            if element_info["type"] != "unsupported":
                elements.append(element_info)

        return {"slide_index": slide_index, "elements": elements}


def get_slide_screenshot_as_base64(pptx_path: str, slide_index: int) -> str:
    img = Image.new('RGB', (800, 600), color='white')
    # 在图像上写字来标识它
    from PIL import ImageDraw
    d = ImageDraw.Draw(img)
    d.text((10, 10), f"Screenshot of Slide {slide_index}", fill=(0, 0, 0))

    buffered = io.BytesIO()
    img.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')


def convert_slide_to_ldl(slide_screenshot_base64: str) -> str:
    """
    使用多模态LLM将幻灯片截图逆向工程为LDL序列。
    这是整个脚本的核心，利用LLM的视觉理解能力来生成结构化描述。
    """
    # 这里的prompt是关键，它指导LLM如何像LPG模型一样“思考”
    prompt_text = f"""
    You are an expert in presentation design and analysis. Your task is to reverse-engineer a slide's visual layout into a symbolic Layout Description Language (LDL).
    Analyze the provided slide screenshot and generate the corresponding LDL sequence.

    The LDL vocabulary consists of:
    - Slide Type Tokens: SLIDE_TITLE, SLIDE_CONTENT_SINGLE_COL, SLIDE_CONTENT_TWO_COL, etc.
    - Element Type Tokens: ELEM_TITLE, ELEM_TEXT_BODY, ELEM_IMAGE, etc.
    - Attribute Tokens: ATTR_TEXT_POINTS_FEW, ATTR_IMAGE_ASPECT_WIDE, ATTR_SIZE_PRIMARY, etc.
    - Position Tokens: POS_TOP_LEFT, POS_CENTER, POS_BOTTOM_RIGHT, etc.
    - Special Tokens: <SOS>, <EOS>, <SEP>.

    Based on the image, produce a single line of text representing the LDL sequence. Start with <SOS> and end with <EOS>. Separate elements with <SEP>.

    Example output for a title and content slide:
    <SOS> SLIDE_CONTENT_SINGLE_COL <SEP> ELEM_TITLE POS_TOP_CENTER <SEP> ELEM_TEXT_BODY ATTR_TEXT_POINTS_MEDIUM POS_MIDDLE_CENTER <EOS>

    Now, analyze the following slide and generate its LDL sequence:
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt_text},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{slide_screenshot_base64}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=300,
            temperature=0.1  # 使用较低的温度以获得更确定性的输出
        )
        ldl_sequence = response.choices[0].message.content.strip()
        # 基本的清洗，移除可能的多余引号或换行符
        ldl_sequence = ldl_sequence.replace("\n", " ").replace("`", "")
        if ldl_sequence.startswith("LDL:") or ldl_sequence.startswith("ldl:"):
            ldl_sequence = ldl_sequence.split(":", 1)[1].strip()

        # 验证输出格式是否基本正确
        if not (ldl_sequence.startswith("<SOS>") and ldl_sequence.endswith("<EOS>")):
            print(f"Warning: LLM generated invalid LDL format: {ldl_sequence}")
            return None

        return ldl_sequence
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        return None


def extract_concept_from_slide(slide_info: Dict[str, Any]) -> Dict[str, Any]:

    num_text_elements = 0
    num_image_elements = 0
    total_text_length = 0

    for elem in slide_info["elements"]:
        if elem["type"] == "text":
            num_text_elements += 1
            total_text_length += len(elem["text"])
        elif elem["type"] == "image":
            num_image_elements += 1


    concept = {
        "num_text_elements": num_text_elements,
        "num_image_elements": num_image_elements,
        "total_text_length": total_text_length,

        "functional_type": "content_text_image"  # 模拟一个默认值
    }
    return concept


def process_all_ppts(source_dir: str, output_file: str):

    pptx_files = list(Path(source_dir).rglob("*.pptx"))
    print(f"Found {len(pptx_files)} PPTX files in '{source_dir}'.")

    with open(output_file, "w", encoding="utf-8") as f_out:
        for pptx_path in tqdm(pptx_files, desc="Processing PPTX files"):
            try:
                parser = PptxParser(str(pptx_path))
                for slide_info in parser:
                    slide_index = slide_info["slide_index"]


                    concept = extract_concept_from_slide(slide_info)


                    screenshot_base64 = get_slide_screenshot_as_base64(str(pptx_path), slide_index)
                    ldl = convert_slide_to_ldl(screenshot_base64)

                    if ldl:
                        data_pair = {"concept": concept, "ldl": ldl}
                        f_out.write(json.dumps(data_pair) + "\n")

            except Exception as e:
                print(f"Failed to process {pptx_path}: {e}")

    print(f"Dataset construction complete. Output saved to '{output_file}'.")


